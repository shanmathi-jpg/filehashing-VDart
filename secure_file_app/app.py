import os
import io
from datetime import datetime
from functools import wraps

from flask import Flask, request, send_file, redirect, url_for, session, render_template
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker, joinedload

from db_test import Base, User, EncryptedFile
from crypto_utils import encrypt_bytes, decrypt_bytes, hash_password, verify_password
from PyPDF2 import PdfReader
from pptx import Presentation

# Flask App
app = Flask(__name__)
app.secret_key = 'change_this_to_a_strong_secret_key'

# Database Setup
db_path = os.path.abspath("files.db")
engine = create_engine(f"sqlite:///{db_path}", echo=False)
SessionLocal = sessionmaker(bind=engine)

# -------------------- Decorators --------------------
def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated

def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('is_admin'):
            return "Access Denied", 403
        return f(*args, **kwargs)
    return decorated

# -------------------- Text Extraction --------------------
def extract_text_from_file(file_bytes, filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".txt":
        return file_bytes.decode("utf-8", errors="replace")
    elif ext == ".pdf":
        reader = PdfReader(io.BytesIO(file_bytes))
        return ''.join(page.extract_text() or '' for page in reader.pages)
    elif ext == ".pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    return ''

# -------------------- Auth Routes --------------------
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        db = SessionLocal()
        username = request.form['username']
        password = request.form['password']
        if db.query(User).filter_by(username=username).first():
            db.close()
            return "Username already exists. <a href='/register'>Try again</a>"
        user = User(username=username, password=hash_password(password))
        db.add(user)
        db.commit()
        db.close()
        return redirect(url_for('login'))
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        db = SessionLocal()
        username = request.form['username']
        password = request.form['password']
        user = db.query(User).filter_by(username=username).first()
        if user and verify_password(password, user.password):
            is_admin = user.is_admin  # ✅ Fix: Cache before session close
            session['user_id'] = user.id
            session['username'] = user.username
            session['is_admin'] = is_admin
            user.last_login = datetime.utcnow()
            db.commit()
            db.close()
            return redirect(url_for('admin_dashboard' if is_admin else 'home'))
        db.close()
        return "Invalid credentials. <a href='/login'>Try again</a>"
    return render_template("login.html", pre_username=request.args.get("username", ""))

@app.route('/logout')
def logout():
    db = SessionLocal()
    user = db.query(User).filter_by(id=session.get('user_id')).first()
    if user:
        user.last_logout = datetime.utcnow()
        db.commit()
    db.close()
    session.clear()
    return redirect(url_for('login'))

# -------------------- User Routes --------------------
@app.route('/')
@login_required
def home():
    return render_template("upload.html", username=session['username'])

@app.route('/upload', methods=['POST'])
@login_required
def upload():
    uploaded_file = request.files.get('file')
    if not uploaded_file:
        return "No file uploaded", 400
    text = extract_text_from_file(uploaded_file.read(), uploaded_file.filename)
    encrypted = encrypt_bytes(text.encode('utf-8'))
    db = SessionLocal()
    file = EncryptedFile(
        filename=uploaded_file.filename,
        data=encrypted,
        user_id=session['user_id']
    )
    db.add(file)
    db.commit()
    db.close()
    return redirect(url_for('list_files'))

@app.route('/files')
@login_required
def list_files():
    db = SessionLocal()
    files = db.query(EncryptedFile).filter_by(user_id=session['user_id']).all()
    db.close()
    return render_template("files.html", files=files, username=session['username'])

@app.route('/view/<int:file_id>')
@login_required
def view_file(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).filter_by(id=file_id, user_id=session['user_id']).first()
    db.close()
    if not file:
        return "Access Denied", 403
    content = decrypt_bytes(file.data).decode("utf-8", errors="replace").strip()
    return render_template("view.html", filename=file.filename, content=content or "Empty file")

@app.route('/download/<int:file_id>')
@login_required
def download(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).filter_by(id=file_id, user_id=session['user_id']).first()
    db.close()
    if not file:
        return "Access Denied", 403
    return send_file(io.BytesIO(decrypt_bytes(file.data)), as_attachment=True, download_name=file.filename)

@app.route('/delete/<int:file_id>')
@login_required
def delete_file(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).filter_by(id=file_id, user_id=session['user_id']).first()
    if file:
        db.delete(file)
        db.commit()
    db.close()
    return redirect(url_for('list_files'))

# -------------------- Admin Routes --------------------
@app.route('/admin')
@admin_required
def admin_dashboard():
    db = SessionLocal()
    files = db.query(EncryptedFile).options(joinedload(EncryptedFile.owner)).all()
    users = db.query(User).all()
    db.close()
    return render_template("admin_files.html", files=files, users=users, username=session["username"])

@app.route('/admin/view/<int:file_id>')
@admin_required
def admin_view_file(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).options(joinedload(EncryptedFile.owner)).filter_by(id=file_id).first()
    db.close()
    if not file:
        return "File not found", 404
    content = decrypt_bytes(file.data).decode("utf-8", errors="replace").strip()
    return render_template("view.html", filename=f"{file.owner.username} - {file.filename}", content=content)

@app.route('/admin/download/<int:file_id>')
@admin_required
def admin_download(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).filter_by(id=file_id).first()
    db.close()
    if not file:
        return "File not found", 404
    return send_file(io.BytesIO(decrypt_bytes(file.data)), as_attachment=True, download_name=file.filename)

@app.route('/admin/delete/<int:file_id>')
@admin_required
def admin_delete_file(file_id):
    db = SessionLocal()
    file = db.query(EncryptedFile).filter_by(id=file_id).first()
    if file:
        db.delete(file)
        db.commit()
    db.close()
    return redirect(url_for("admin_dashboard"))

# -------------------- Run App --------------------
if __name__ == "__main__":
    app.run(debug=True)
